from __future__ import annotations

import re
from dataclasses import dataclass

from pptx import Presentation

_NO_CHECKPOINT = "[NO-CHECKPOINT]"
_SECTION_DELIMITER = "---"


@dataclass(frozen=True)
class Section:
    index: int  # global, 0-based — the orchestrator's position
    slide_index: int  # 1-based, matches PowerPoint
    section_index: int  # 0-based within the slide
    text: str
    checkpoint: bool


@dataclass(frozen=True)
class LectureScript:
    sections: list[Section]  # flat, ordered; the orchestrator's playlist
    slide_count: int
    full_narration: str  # all section text joined — used in §7.3 grounding


@dataclass(frozen=True)
class ValidationResult:
    errors: list[str]
    warnings: list[str]

    @property
    def is_valid(self) -> bool:
        return not self.errors


_MAX_SECTION_CHARS = 1500


def validate(pptx_path: str) -> ValidationResult:
    """Deck-level checks from §5.4: file opens, has slides, every slide is
    narrated (fatal), no section over the TTS-latency guard (warning).

    The other §5.4 checks (qa_material.md, Ollama, Piper, ffmpeg) belong to
    the modules that own those concerns, not here.
    """
    try:
        prs = Presentation(pptx_path)
    except Exception as e:
        return ValidationResult(errors=[f"Could not open '{pptx_path}': {e}"], warnings=[])

    if len(prs.slides) < 1:
        return ValidationResult(errors=["Deck has no slides."], warnings=[])

    script = parse(pptx_path)
    errors: list[str] = []
    warnings: list[str] = []

    narrated_slides = {s.slide_index for s in script.sections}
    for slide_num in range(1, script.slide_count + 1):
        if slide_num not in narrated_slides:
            errors.append(
                f"Slide {slide_num} has no speaker notes — every slide must be narrated."
            )

    for s in script.sections:
        if len(s.text) > _MAX_SECTION_CHARS:
            warnings.append(
                f"Slide {s.slide_index} section {s.section_index} is {len(s.text)} "
                f"chars, over the {_MAX_SECTION_CHARS}-char TTS-latency guard."
            )

    return ValidationResult(errors=errors, warnings=warnings)


def parse(pptx_path: str) -> LectureScript:
    prs = Presentation(pptx_path)
    sections: list[Section] = []
    global_index = 0

    for slide_num, slide in enumerate(prs.slides, start=1):
        notes_text = ""
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text

        section_idx = 0
        for raw in _split_sections(notes_text):
            has_marker = _NO_CHECKPOINT in raw
            text = _clean(raw.replace(_NO_CHECKPOINT, "") if has_marker else raw)
            if not text:
                continue  # blank sections discarded

            sections.append(
                Section(
                    index=global_index,
                    slide_index=slide_num,
                    section_index=section_idx,
                    text=text,
                    checkpoint=not has_marker,
                )
            )
            global_index += 1
            section_idx += 1

    # Final section always has a checkpoint regardless of marker — the
    # queue must drain before FINISHED.
    if sections and not sections[-1].checkpoint:
        last = sections[-1]
        sections[-1] = Section(
            index=last.index,
            slide_index=last.slide_index,
            section_index=last.section_index,
            text=last.text,
            checkpoint=True,
        )

    full_narration = " ".join(s.text for s in sections)

    return LectureScript(
        sections=sections,
        slide_count=len(prs.slides),
        full_narration=full_narration,
    )


def _split_sections(notes_text: str) -> list[str]:
    # Split only on a line that is *exactly* "---" after stripping — a
    # naive substring split would also cut inside narration text that
    # happens to contain "---".
    lines = notes_text.split("\n")
    sections: list[str] = []
    current: list[str] = []
    for line in lines:
        if line.strip() == _SECTION_DELIMITER:
            sections.append("\n".join(current))
            current = []
        else:
            current.append(line)
    sections.append("\n".join(current))
    return sections


def _clean(text: str) -> str:
    return re.sub(r"\s+", " ", text).strip()
