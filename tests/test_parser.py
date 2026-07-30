from pathlib import Path

from pptx import Presentation

from app.script_parser import parse


def _make_deck(tmp_path: Path, notes_per_slide: list[str]) -> str:
    prs = Presentation()
    layout = prs.slide_layouts[6]  # blank layout
    for notes in notes_per_slide:
        slide = prs.slides.add_slide(layout)
        slide.notes_slide.notes_text_frame.text = notes
    path = tmp_path / "test_deck.pptx"
    prs.save(str(path))
    return str(path)


def test_single_section_per_slide(tmp_path):
    deck = _make_deck(tmp_path, ["First slide narration.", "Second slide narration."])
    script = parse(deck)

    assert script.slide_count == 2
    assert [s.text for s in script.sections] == [
        "First slide narration.",
        "Second slide narration.",
    ]
    assert [s.checkpoint for s in script.sections] == [True, True]
    assert [s.slide_index for s in script.sections] == [1, 2]
    assert [s.section_index for s in script.sections] == [0, 0]
    assert [s.index for s in script.sections] == [0, 1]


def test_multi_section_slide_splits_on_bare_dashes(tmp_path):
    notes = "Part one.\n---\nPart two.\n---\nPart three."
    deck = _make_deck(tmp_path, [notes])
    script = parse(deck)

    assert [s.text for s in script.sections] == ["Part one.", "Part two.", "Part three."]
    assert [s.section_index for s in script.sections] == [0, 1, 2]
    assert [s.slide_index for s in script.sections] == [1, 1, 1]


def test_dashes_inside_a_line_do_not_split():
    # only a line that is *exactly* "---" is a delimiter
    from app.script_parser import _split_sections

    result = _split_sections("Score was 10---9 in the game.")
    assert result == ["Score was 10---9 in the game."]


def test_no_checkpoint_suppresses_checkpoint_after_that_section(tmp_path):
    notes = "Section A.\n---\n[NO-CHECKPOINT]\nSection B.\n---\nSection C."
    deck = _make_deck(tmp_path, [notes])
    script = parse(deck)

    assert [s.text for s in script.sections] == ["Section A.", "Section B.", "Section C."]
    # A has a checkpoint (default); B suppresses the checkpoint *after* B;
    # C is the final section, so it gets a checkpoint regardless.
    assert [s.checkpoint for s in script.sections] == [True, False, True]


def test_marker_stripped_from_text(tmp_path):
    deck = _make_deck(tmp_path, ["[NO-CHECKPOINT]\nText after the marker."])
    script = parse(deck)

    assert script.sections[0].text == "Text after the marker."
    assert "[NO-CHECKPOINT]" not in script.sections[0].text


def test_final_section_always_has_checkpoint_even_if_marked(tmp_path):
    deck = _make_deck(tmp_path, ["[NO-CHECKPOINT]\nOnly section, marked anyway."])
    script = parse(deck)

    assert len(script.sections) == 1
    assert script.sections[0].checkpoint is True


def test_blank_sections_are_discarded(tmp_path):
    notes = "Real content.\n---\n   \n---\nMore real content."
    deck = _make_deck(tmp_path, [notes])
    script = parse(deck)

    assert [s.text for s in script.sections] == ["Real content.", "More real content."]
    assert [s.section_index for s in script.sections] == [0, 1]


def test_slide_with_no_notes_produces_no_sections(tmp_path):
    deck = _make_deck(tmp_path, ["Has notes.", ""])
    script = parse(deck)

    assert script.slide_count == 2
    assert len(script.sections) == 1
    assert script.sections[0].slide_index == 1


def test_full_narration_joins_all_section_text(tmp_path):
    deck = _make_deck(tmp_path, ["First slide.", "Second slide, part one.\n---\nSecond slide, part two."])
    script = parse(deck)

    assert script.full_narration == "First slide. Second slide, part one. Second slide, part two."


def test_global_index_is_sequential_across_slides(tmp_path):
    deck = _make_deck(
        tmp_path,
        ["A.\n---\nB.", "C."],
    )
    script = parse(deck)

    assert [s.index for s in script.sections] == [0, 1, 2]
    assert [(s.slide_index, s.section_index) for s in script.sections] == [
        (1, 0),
        (1, 1),
        (2, 0),
    ]


def test_whitespace_is_trimmed_and_internal_newlines_collapsed(tmp_path):
    deck = _make_deck(tmp_path, ["  Line one.\n  Line two continues.  "])
    script = parse(deck)

    assert script.sections[0].text == "Line one. Line two continues."
